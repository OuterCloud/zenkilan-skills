#!/usr/bin/env python3
"""
fill_template.py - Fill a .pptx template with structured content data.

Usage:
    python fill_template.py --template <template.pptx> --data <content.json> --output <output.pptx>

Data JSON format:
{
  "slides": [
    {
      "layout": "Title Slide",          // Use a named layout
      "placeholders": {
        "0": {"type": "text", "content": "Title text"},
        "1": {"type": "text", "content": "Subtitle"}
      }
    },
    {
      "duplicate_from": 0,              // Duplicate existing slide at index 0
      "placeholders": {
        "0": {"type": "text", "content": "New title"}
      }
    },
    {
      "layout": "Two Content",
      "placeholders": {
        "1": {"type": "table", "content": {"headers": ["A","B"], "rows": [["1","2"]]}},
        "2": {"type": "image", "content": "/path/to/image.png"}
      }
    }
  ]
}

Rich text format for "text" type:
- Plain string: used as-is, preserving template formatting
- Array of paragraph objects for rich text:
  [
    {"text": "Bold title", "bold": true, "size": 28},
    {"text": "Normal paragraph"},
    {"text": "Red text", "color": "FF0000", "bold": true, "size": 14}
  ]
"""

import sys
import json
import copy
import argparse
from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def duplicate_slide(prs, slide_index):
    """
    Duplicate an existing slide by copying its XML and relationships.
    python-pptx doesn't support this natively, so we manipulate the XML directly.
    
    Returns the new slide object.
    """
    source_slide = prs.slides[slide_index]
    slide_layout = source_slide.slide_layout

    # Add a new slide with the same layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy the source slide's XML tree (deep copy)
    # We need to copy the spTree (shape tree) from source to new slide
    # First, remove existing shapes from new slide (except the background)
    new_slide_elem = new_slide._element
    source_slide_elem = source_slide._element

    # Remove all child elements from new slide's spTree
    new_sp_tree = new_slide_elem.find(
        './/{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}spTree'
    )
    if new_sp_tree is None:
        new_sp_tree = new_slide_elem.find(
            './/{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
        )

    # Deep copy the entire cSld element from source
    source_cSld = source_slide_elem.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
    )
    new_cSld = new_slide_elem.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
    )

    if source_cSld is not None and new_cSld is not None:
        # Replace the new slide's cSld with a deep copy of source
        new_slide_elem.replace(new_cSld, copy.deepcopy(source_cSld))

    # Copy relationships (images, charts, etc.)
    for rel in source_slide.part.rels.values():
        # Skip layout and notesSlide rels
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            # For embedded parts (images, etc.), copy the target part
            try:
                new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
            except Exception:
                # Some rels may not be copyable; skip gracefully
                pass

    return new_slide


def find_layout_by_name(prs, layout_name):
    """Find a slide layout by name (case-insensitive partial match)."""
    # Exact match first
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return layout
    # Case-insensitive match
    for layout in prs.slide_layouts:
        if layout.name.lower() == layout_name.lower():
            return layout
    # Partial match
    for layout in prs.slide_layouts:
        if layout_name.lower() in layout.name.lower():
            return layout
    return None


def apply_text_content(placeholder, content):
    """
    Apply text content to a placeholder.
    
    content can be:
    - A plain string (simple replacement, preserves first paragraph formatting)
    - A list of paragraph objects with optional formatting:
      [{"text": "...", "bold": true, "size": 28, "color": "FF0000", "align": "center"}]
    """
    if not placeholder.has_text_frame:
        return

    tf = placeholder.text_frame

    if isinstance(content, str):
        # Simple text: split by newlines, preserve first paragraph's formatting
        lines = content.split("\n")

        # Use first existing paragraph for first line
        if tf.paragraphs:
            first_para = tf.paragraphs[0]
            # Preserve existing run formatting if possible
            if first_para.runs:
                first_para.runs[0].text = lines[0]
                # Remove extra runs
                for run in first_para.runs[1:]:
                    first_para._p.remove(run._r)
            else:
                first_para.text = lines[0]
        else:
            tf.text = lines[0]

        # Add remaining lines as new paragraphs
        for line in lines[1:]:
            para = tf.add_paragraph()
            # Copy formatting from first paragraph if available
            if tf.paragraphs[0].runs:
                run = para.add_run()
                run.text = line
                src_run = tf.paragraphs[0].runs[0]
                if src_run.font.size:
                    run.font.size = src_run.font.size
                try:
                    if src_run.font.color and src_run.font.color.type is not None:
                        run.font.color.rgb = src_run.font.color.rgb
                except (AttributeError, TypeError):
                    pass
            else:
                para.text = line

    elif isinstance(content, list):
        # Rich text: list of paragraph objects
        # Clear existing paragraphs (keep the first one for formatting reference)
        tf.clear()

        for i, para_data in enumerate(content):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            text = para_data.get("text", "")
            run = para.add_run()
            run.text = text

            # Apply formatting
            if "bold" in para_data:
                run.font.bold = para_data["bold"]
            if "italic" in para_data:
                run.font.italic = para_data["italic"]
            if "size" in para_data:
                run.font.size = Pt(para_data["size"])
            if "color" in para_data:
                run.font.color.rgb = RGBColor.from_string(para_data["color"])
            if "font_name" in para_data:
                run.font.name = para_data["font_name"]

            # Paragraph-level formatting
            if "align" in para_data:
                align_map = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT,
                    "justify": PP_ALIGN.JUSTIFY,
                }
                para.alignment = align_map.get(para_data["align"], PP_ALIGN.LEFT)

            if "space_before" in para_data:
                para.space_before = Pt(para_data["space_before"])
            if "space_after" in para_data:
                para.space_after = Pt(para_data["space_after"])


def apply_table_content(placeholder, content):
    """
    Apply table content to a placeholder.
    content: {"headers": ["A", "B"], "rows": [["1","2"], ["3","4"]]}
    """
    headers = content.get("headers", [])
    rows = content.get("rows", [])
    col_count = len(headers)
    row_count = len(rows) + 1  # +1 for header row

    if col_count == 0:
        return

    # Get placeholder dimensions
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    # We need to add a table to the slide (not into the placeholder directly)
    slide = placeholder._element.getparent().getparent()
    # Find the slide object - traverse up to get the actual slide
    # Instead, we'll use the graphic frame approach

    # python-pptx supports inserting tables via placeholder if it's a table placeholder
    # For non-table placeholders, we insert a table using the placeholder's position
    try:
        # Try using placeholder's insert_table (works for table/object placeholders)
        table_frame = placeholder.insert_table(row_count, col_count)
        table = table_frame.table
    except (AttributeError, TypeError):
        # Fallback: can't insert table into this placeholder type
        # The table will need to be added at slide level
        return

    # Fill header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = str(header)
        # Make header bold
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.bold = True

    # Fill data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_value in enumerate(row_data):
            if col_idx < col_count:
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_value)


def apply_image_content(placeholder, image_path):
    """Apply an image to a placeholder."""
    image_path = str(image_path)
    if not Path(image_path).exists():
        print(f"  Warning: Image not found: {image_path}", file=sys.stderr)
        return

    try:
        # Use insert_picture for picture placeholders
        placeholder.insert_picture(image_path)
    except (AttributeError, TypeError):
        print(f"  Warning: Cannot insert image into placeholder {placeholder.name}", file=sys.stderr)


def fill_placeholder(placeholder, ph_data):
    """Fill a single placeholder based on its type specification."""
    ph_type = ph_data.get("type", "text")
    content = ph_data.get("content", "")

    if ph_type == "text":
        apply_text_content(placeholder, content)
    elif ph_type == "table":
        apply_table_content(placeholder, content)
    elif ph_type == "image":
        apply_image_content(placeholder, content)
    else:
        print(f"  Warning: Unknown placeholder type '{ph_type}'", file=sys.stderr)


def process_slide(prs, slide_spec, created_slides):
    """
    Process a single slide specification.
    Returns the created/modified slide.
    """
    if "duplicate_from" in slide_spec:
        # Duplicate an existing slide
        source_idx = slide_spec["duplicate_from"]
        if source_idx < 0 or source_idx >= len(prs.slides):
            print(f"  Error: duplicate_from index {source_idx} out of range", file=sys.stderr)
            return None
        slide = duplicate_slide(prs, source_idx)
    elif "layout" in slide_spec:
        # Create new slide from layout
        layout_name = slide_spec["layout"]
        layout = find_layout_by_name(prs, layout_name)
        if layout is None:
            print(f"  Error: Layout '{layout_name}' not found. Available layouts:", file=sys.stderr)
            for l in prs.slide_layouts:
                print(f"    - {l.name}", file=sys.stderr)
            return None
        slide = prs.slides.add_slide(layout)
    else:
        print("  Error: Slide spec must have 'layout' or 'duplicate_from'", file=sys.stderr)
        return None

    # Fill placeholders
    placeholders_data = slide_spec.get("placeholders", {})
    for idx_str, ph_data in placeholders_data.items():
        try:
            idx = int(idx_str)
        except ValueError:
            print(f"  Warning: Invalid placeholder index '{idx_str}'", file=sys.stderr)
            continue

        # Find placeholder by idx
        placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                placeholder = ph
                break

        if placeholder is None:
            print(f"  Warning: Placeholder idx={idx} not found on slide", file=sys.stderr)
            continue

        fill_placeholder(placeholder, ph_data)

    return slide


def main():
    parser = argparse.ArgumentParser(
        description="Fill a .pptx template with structured content data."
    )
    parser.add_argument("--template", required=True, help="Path to template .pptx file")
    parser.add_argument("--data", required=True, help="Path to content JSON file")
    parser.add_argument("--output", required=True, help="Output .pptx file path")
    parser.add_argument(
        "--keep-template-slides",
        action="store_true",
        help="Keep original template slides (by default they are removed)",
    )
    args = parser.parse_args()

    # Load template
    try:
        prs = Presentation(args.template)
    except Exception as e:
        print(f"Error loading template: {e}", file=sys.stderr)
        sys.exit(1)

    # Load data
    try:
        with open(args.data, "r", encoding="utf-8") as f:
            data = json.load(f)
    except Exception as e:
        print(f"Error loading data JSON: {e}", file=sys.stderr)
        sys.exit(1)

    slides_data = data.get("slides", [])
    if not slides_data:
        print("Warning: No slides defined in data JSON", file=sys.stderr)

    # Record original slide count for potential removal
    original_slide_count = len(prs.slides)

    # Process each slide spec
    created_slides = []
    for i, slide_spec in enumerate(slides_data):
        print(f"Processing slide {i + 1}/{len(slides_data)}...", file=sys.stderr)
        slide = process_slide(prs, slide_spec, created_slides)
        if slide:
            created_slides.append(slide)

    # Remove original template slides if not keeping them
    if not args.keep_template_slides and original_slide_count > 0:
        # Remove slides from the beginning (original template slides)
        # We need to remove by manipulating the XML directly
        slide_list = prs.slides._sldIdLst
        slide_ids = list(slide_list)
        for i in range(original_slide_count):
            if i < len(slide_ids):
                slide_id_elem = slide_ids[i]
                slide_list.remove(slide_id_elem)

    # Save output
    try:
        output_path = Path(args.output)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        print(f"Output saved to: {args.output}", file=sys.stderr)
    except Exception as e:
        print(f"Error saving output: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
