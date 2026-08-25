#!/usr/bin/env python3
"""
analyze_template.py - Analyze a .pptx template and output its structure as JSON.

Usage:
    python analyze_template.py <template.pptx>

Output: JSON describing all slide layouts, their placeholders, and existing slides.
"""

import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


def emu_to_inches(emu_value):
    """Convert EMU to inches, rounded to 2 decimal places."""
    if emu_value is None:
        return None
    return round(emu_value / 914400, 2)


def placeholder_type_name(ph_type):
    """Get human-readable placeholder type name."""
    type_map = {
        PP_PLACEHOLDER.TITLE: "title",
        PP_PLACEHOLDER.CENTER_TITLE: "center_title",
        PP_PLACEHOLDER.SUBTITLE: "subtitle",
        PP_PLACEHOLDER.BODY: "body",
        PP_PLACEHOLDER.OBJECT: "object",
        PP_PLACEHOLDER.DATE: "date",
        PP_PLACEHOLDER.FOOTER: "footer",
        PP_PLACEHOLDER.SLIDE_NUMBER: "slide_number",
        PP_PLACEHOLDER.TABLE: "table",
        PP_PLACEHOLDER.CHART: "chart",
        PP_PLACEHOLDER.PICTURE: "picture",
        PP_PLACEHOLDER.BITMAP: "bitmap",
        PP_PLACEHOLDER.MEDIA_CLIP: "media_clip",
        PP_PLACEHOLDER.ORG_CHART: "org_chart",
    }
    return type_map.get(ph_type, str(ph_type) if ph_type else "unknown")


def get_text_summary(text_frame, max_len=80):
    """Get a summary of text content from a text frame."""
    if text_frame is None:
        return None
    full_text = text_frame.text.strip()
    if not full_text:
        return ""
    if len(full_text) <= max_len:
        return full_text
    return full_text[:max_len] + "..."


def analyze_placeholder(placeholder):
    """Analyze a single placeholder shape."""
    info = {
        "idx": placeholder.placeholder_format.idx,
        "name": placeholder.name,
        "type": placeholder_type_name(placeholder.placeholder_format.type),
        "position": {
            "left_inches": emu_to_inches(placeholder.left),
            "top_inches": emu_to_inches(placeholder.top),
            "width_inches": emu_to_inches(placeholder.width),
            "height_inches": emu_to_inches(placeholder.height),
        },
    }
    # Add content summary if it has text
    if placeholder.has_text_frame:
        summary = get_text_summary(placeholder.text_frame)
        if summary is not None:
            info["content_summary"] = summary
    return info


def analyze_layout(layout):
    """Analyze a single slide layout."""
    placeholders = []
    for ph in layout.placeholders:
        placeholders.append(analyze_placeholder(ph))
    # Sort by idx for consistent output
    placeholders.sort(key=lambda x: x["idx"])
    return {
        "name": layout.name,
        "placeholders": placeholders,
    }


def analyze_slide(slide, slide_index):
    """Analyze a single slide."""
    # Determine layout name
    layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"

    placeholders = []
    for ph in slide.placeholders:
        ph_info = analyze_placeholder(ph)
        # For existing slides, always try to get content
        if ph.has_text_frame:
            ph_info["content_summary"] = get_text_summary(ph.text_frame, max_len=120)
        elif ph.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ph_info["content_summary"] = "[image]"
        elif ph.shape_type == MSO_SHAPE_TYPE.TABLE:
            ph_info["content_summary"] = "[table]"
        placeholders.append(ph_info)

    placeholders.sort(key=lambda x: x["idx"])

    return {
        "index": slide_index,
        "layout": layout_name,
        "placeholders": placeholders,
    }


def analyze_template(pptx_path):
    """Analyze a .pptx template file and return structured data."""
    prs = Presentation(pptx_path)

    result = {
        "file": pptx_path,
        "slide_width_inches": emu_to_inches(prs.slide_width),
        "slide_height_inches": emu_to_inches(prs.slide_height),
        "layouts": [],
        "existing_slides": [],
    }

    # Analyze all slide layouts
    for layout in prs.slide_layouts:
        result["layouts"].append(analyze_layout(layout))

    # Analyze existing slides
    for i, slide in enumerate(prs.slides):
        result["existing_slides"].append(analyze_slide(slide, i))

    return result


def main():
    if len(sys.argv) != 2:
        print("Usage: python analyze_template.py <template.pptx>", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]

    try:
        result = analyze_template(pptx_path)
        print(json.dumps(result, ensure_ascii=False, indent=2))
    except Exception as e:
        print(json.dumps({"error": str(e)}), file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
